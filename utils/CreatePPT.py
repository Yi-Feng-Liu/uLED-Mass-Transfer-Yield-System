import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.dml import MSO_THEME_COLOR
import os
from glob import glob
from datetime import datetime


class PPTmain():
    """
    Create PPT file by image name, if save path did not exit ppt file that will be write to log file to record.
    """
    def __init__(self):
        self.prs = Presentation('./utils/TEST.pptx')
        self.byPer24Hour = glob('./report_production/*Per24Hour.jpg')
        self.removeDefectImage = glob('./report_production/*rmDefect.png')
        self.originalDefectImage = glob('./report_production/*original.png')
        self.imgls = glob('./report_production/*.bmp')
        self.byday = glob('./report_production/*_byDay.jpg')
        
    def addSlide(self, titleText:str, image_path=None, left=Inches(0.6), top=Inches(1.5), width=None,   
                 height=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.placeholders[0].text = titleText
        slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(30)
        slide.placeholders[0].text_frame.paragraphs[0].font.bold = True
        slide.placeholders[0].text_frame.paragraphs[0].font.name = 'Time New Roman'
        slide.placeholders[0].text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        slide.shapes.add_picture(image_path, left=left, top=top, width=width, height=height)
        return slide

    def dailyReport(self):
        """Create PPT file for table images of data from last 24 hours. 
        """
        if len(self.byPer24Hour) != 0:
            for i in range(len(self.byPer24Hour)):
                img = self.byPer24Hour[i]
                if self.byPer24Hour[i].endswith("Per24Hour.jpg"):
                    OPID = self.byPer24Hour[i].split('_')[1].split('\\')[-1]
                    self.addSlide(f"Daily Report_{OPID} Analysis", image_path=img, width=Inches(12), height=Inches(6))

            # print(glob('.\\report_production\\*'))
            sheetIDls = [i.split('\\')[-1].split('_')[0] for i in glob('.\\report_production\\*png')]
            # print(set(sheetIDls))
            all_file = glob('.\\report_production\\*')
            all_file = sorted(all_file, reverse=True)

            for sheetID in set(sheetIDls):
                for files in all_file:
                    file_sheet_ID = files.split('.')[-2].split('\\')[-1].split("_")[0]
                    # order by sheet id
                    if sheetID == file_sheet_ID:
                        fileName = files.split('\\')[-1].split('.')[-2]
                        self.addSlide(f"{fileName}", image_path=files, width=Inches(12), height=Inches(6))
                    else:
                        continue
        
            pptx_name = f'Daily_report_{datetime.today().weekday()}.pptx'
            self.prs.save(f"./report_production/{pptx_name}")
            # os.startfile(".\\report_production\\Daily_report.pptx")
            return pptx_name

    def weeklyReport(self):
        """Create PPT file for table images of data from two weeks ago.  
        """
        if len(self.byday) != 0:
            for i in range(len(self.byday)):
                img = self.byday[i]
                if self.byday[i].endswith("byDay.jpg"):
                    OPID = self.byday[i].split('_')[1].split('\\')[-1]
                    self.addSlide(f"Weekly Report_{OPID} Analysis", image_path=img, width=Inches(12), height=Inches(6))

            sheetIDls = [i.split('\\')[-1].split('_')[0] for i in glob('.\\report_production\\*png')]
            # print(set(sheetIDls))
            all_file = glob('.\\report_production\\*')
            all_file = sorted(all_file, reverse=True)

            for sheetID in set(sheetIDls):
                for files in all_file:
                    file_sheet_ID = files.split('.')[-2].split('\\')[-1].split("_")[0]
                    # order by sheet id
                    if sheetID == file_sheet_ID:
                        fileName = files.split('\\')[-1].split('.')[-2]
                        self.addSlide(f"{fileName}", image_path=files, width=Inches(12), height=Inches(6))
                    else:
                        continue
            pptx_name = f'Weekly_report_{datetime.today().weekday()}.pptx'
            self.prs.save(f"./report_production/{pptx_name}")
            # os.startfile(f".\\report_production\\{pptx_name}")
            return pptx_name

if __name__ == '__main__':

    PPTmain().dailyReport()
